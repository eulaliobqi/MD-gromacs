"""
Atualiza apresentacao_md_spodoptera.pptx adicionando slide de resultados BEN.
Insere novo slide antes de Conclusões com dados de docking+DM da benzamidina.
"""
import copy
from pathlib import Path
from lxml import etree

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt

PPTX_IN  = Path(__file__).parent.parent / "apresentacao_md_spodoptera.pptx"
PPTX_OUT = Path(__file__).parent.parent / "apresentacao_md_spodoptera.pptx"

# Cores do design system (inferidas da apresentação existente)
AZUL_ESCURO  = RGBColor(0x1F, 0x49, 0x7D)
AZUL_MEDIO   = RGBColor(0x2E, 0x75, 0xB6)
VERDE        = RGBColor(0x37, 0x86, 0x48)
VERMELHO     = RGBColor(0xC0, 0x00, 0x00)
AMARELO      = RGBColor(0xFF, 0xC0, 0x00)
CINZA        = RGBColor(0xF2, 0xF2, 0xF2)
BRANCO       = RGBColor(0xFF, 0xFF, 0xFF)
LARANJA      = RGBColor(0xED, 0x7D, 0x31)
PRETO        = RGBColor(0x00, 0x00, 0x00)

W  = Inches(13.33)
H  = Inches(7.5)


def add_textbox(slide, left, top, width, height, text, fontsize=18,
                bold=False, color=PRETO, bg=None, align=PP_ALIGN.LEFT,
                italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.line.fill.background()
    if bg:
        txBox.fill.solid()
        txBox.fill.fore_color.rgb = bg
    else:
        txBox.fill.background()
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(fontsize)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_table_slide(prs, insert_idx: int):
    """
    Adiciona slide com tabela comparativa completa GORE4 / SKTI / BEN.
    Inserido na posição insert_idx (0-based).
    """
    blank_layout = prs.slide_layouts[6]  # layout em branco
    slide = prs.slides.add_slide(blank_layout)

    # Mover para posição correta
    xml_slides = prs.slides._sldIdLst
    slides = prs.slides._sldIdLst.findall(
        '{http://schemas.openxmlformats.org/presentationml/2006/main}sldId'
    )
    last = slides[-1]
    xml_slides.remove(last)
    xml_slides.insert(insert_idx, last)

    # Fundo azul escuro
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = AZUL_ESCURO

    # Título
    add_textbox(slide,
                left=Inches(0.3), top=Inches(0.15),
                width=Inches(12.7), height=Inches(0.7),
                text="Benzamidina (BEN) — Controle Positivo Estrutural",
                fontsize=26, bold=True, color=BRANCO,
                align=PP_ALIGN.CENTER)

    # Subtítulo docking
    add_textbox(slide,
                left=Inches(0.3), top=Inches(0.85),
                width=Inches(6.0), height=Inches(0.4),
                text="Docking AutoDock Vina (4 receptores, pH 8,2)",
                fontsize=14, bold=True, color=AMARELO)

    # Tabela docking scores
    dock_data = [
        ["Receptor", "Score Vina (kcal/mol)", "S1", "Hierarquia"],
        ["QCL936",   "−5,733 (melhor)",       "Asp241 ✓",  "1°"],
        ["XP273",    "−5,484",                "Ile229",     "2°"],
        ["XP352",    "−4,975",                "Asp262 ✓",   "3°"],
        ["ACR157",   "−4,953 (pior)",         "Ile205 ✗",   "4°"],
    ]
    col_w = [Inches(1.5), Inches(2.2), Inches(1.5), Inches(1.0)]
    row_h = Inches(0.32)
    tbl_top = Inches(1.3)
    tbl_left = Inches(0.3)

    tbl = slide.shapes.add_table(
        len(dock_data), len(dock_data[0]),
        tbl_left, tbl_top,
        sum(col_w), row_h * len(dock_data)
    ).table

    for c_idx, cw in enumerate(col_w):
        tbl.columns[c_idx].width = cw

    for r, row in enumerate(dock_data):
        for c, val in enumerate(row):
            cell = tbl.cell(r, c)
            cell.text = val
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.CENTER
            run = p.runs[0]
            run.font.size = Pt(11)
            run.font.name = "Calibri"
            if r == 0:
                run.font.bold = True
                run.font.color.rgb = BRANCO
                cell.fill.solid()
                cell.fill.fore_color.rgb = AZUL_MEDIO
            elif "melhor" in val or "1°" in val:
                run.font.bold = True
                run.font.color.rgb = VERDE
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xE2, 0xEF, 0xDA)
            elif "Ile205" in val or "pior" in val or "4°" in val:
                run.font.color.rgb = VERMELHO
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xFF, 0xEB, 0xEB)
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = RGBColor(0xF2, 0xF2, 0xF2)

    # Nota sobre ponte salina
    add_textbox(slide,
                left=Inches(0.3), top=Inches(3.05),
                width=Inches(6.0), height=Inches(0.4),
                text="BEN exige Asp no S1 para ponte salina (amidínio +1 ↔ Asp⁻)",
                fontsize=11, italic=True, color=RGBColor(0xCC, 0xCC, 0x00))

    # Separador vertical
    line = slide.shapes.add_connector(1, Inches(6.6), Inches(0.8), Inches(6.6), Inches(7.1))
    line.line.color.rgb = RGBColor(0x60, 0x90, 0xC0)
    line.line.width = Pt(1)

    # Lado direito — DM ACR157-BEN
    add_textbox(slide,
                left=Inches(6.8), top=Inches(0.85),
                width=Inches(6.3), height=Inches(0.4),
                text="DM ACR157-BEN (200 ns) — Resultado",
                fontsize=14, bold=True, color=AMARELO)

    # Caixa de destaque: DISSOCIAÇÃO
    add_textbox(slide,
                left=Inches(6.8), top=Inches(1.3),
                width=Inches(6.3), height=Inches(0.9),
                text="DISSOCIAÇÃO COMPLETA  ~95 ns",
                fontsize=22, bold=True, color=VERMELHO,
                align=PP_ALIGN.CENTER,
                bg=RGBColor(0x40, 0x10, 0x10))

    # Parâmetros chave
    params = [
        ("RMSD backbone (200 ns)",   "0,146 ± 0,019 nm",   "✅ receptor estável"),
        ("Contatos receptor–BEN",     "68,8 ± 72,4 átomos", "❌ bimodal / dissociação"),
        ("SASA BEN (200 ns)",         "2,74 ± 0,141 nm²",   "—"),
        ("Distância Ile205 (0–95 ns)","~0,35 nm borderline", "⚠️ contacto fraco"),
        ("Distância após ~95 ns",     "3–6 nm (todos res.)", "❌ BEN livre no solvente"),
    ]
    for idx, (label, val, status) in enumerate(params):
        y = Inches(2.35) + Inches(0.52) * idx
        add_textbox(slide,
                    left=Inches(6.8), top=y,
                    width=Inches(3.0), height=Inches(0.45),
                    text=label, fontsize=11, color=RGBColor(0xCC, 0xDD, 0xFF))
        add_textbox(slide,
                    left=Inches(9.8), top=y,
                    width=Inches(1.8), height=Inches(0.45),
                    text=val, fontsize=11, bold=True, color=BRANCO)
        sc = VERMELHO if "❌" in status else (
             AMARELO if "⚠️" in status else VERDE)
        add_textbox(slide,
                    left=Inches(11.6), top=y,
                    width=Inches(1.5), height=Inches(0.45),
                    text=status, fontsize=10, color=sc)

    # Interpretação
    add_textbox(slide,
                left=Inches(6.8), top=Inches(5.05),
                width=Inches(6.3), height=Inches(1.8),
                text=(
                    "Interpretação:\n"
                    "• ACR157 tem Ile205 (hidrofóbico) no S1 → sem ponte salina para amidínio\n"
                    "• BEN perde âncora primária → dissociação irreversível antes de 100 ns\n"
                    "• Consistente com menor score Vina (−4,95 kcal/mol, 4°/4)\n"
                    "• QCL936/XP352 (Asp no S1) → candidatos para BEN ligada — DM pendente"
                ),
                fontsize=11, color=BRANCO)

    # Rodapé
    add_textbox(slide,
                left=Inches(0.3), top=Inches(7.0),
                width=Inches(12.7), height=Inches(0.3),
                text="Campo de força: AMBER99SB-ILDN (receptor) + GAFF2 (BEN, carga +1 amidínio) | 300 K | 0,10 M KCl | pH 8,2",
                fontsize=9, italic=True, color=RGBColor(0x88, 0xAA, 0xCC),
                align=PP_ALIGN.CENTER)

    return slide


def main():
    prs = Presentation(str(PPTX_IN))

    # Encontrar índice do slide "Conclusões" (slide 14 → índice 13)
    conclusions_idx = None
    for i, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if hasattr(shape, 'text') and 'Conclus' in shape.text:
                conclusions_idx = i
                break
        if conclusions_idx is not None:
            break

    if conclusions_idx is None:
        conclusions_idx = len(prs.slides) - 1  # antes do último

    print(f"[INFO] Inserindo slide BEN antes do slide {conclusions_idx + 1} (Conclusões)")

    add_table_slide(prs, insert_idx=conclusions_idx)

    prs.save(str(PPTX_OUT))
    print(f"[OK] Apresentação salva: {PPTX_OUT}")
    print(f"[OK] Total de slides: {len(prs.slides)}")


if __name__ == '__main__':
    main()
