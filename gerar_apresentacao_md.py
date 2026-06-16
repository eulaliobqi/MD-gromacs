#!/usr/bin/env python3
"""
Gera apresentacao_md_spodoptera.pptx — 14 slides
Uso: python gerar_apresentacao_md.py
"""
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

BASE = os.path.dirname(os.path.abspath(__file__))
IMG  = os.path.join(BASE, "results-MD")

# ── Dimensões 16:9 ────────────────────────────────────────────────────────────
W = Inches(13.33)
H = Inches(7.50)

# ── Paleta ────────────────────────────────────────────────────────────────────
DB    = RGBColor(0x1B, 0x3A, 0x6B)  # dark blue
TEAL  = RGBColor(0x0A, 0x7B, 0x8A)  # teal
ORG   = RGBColor(0xD9, 0x72, 0x12)  # orange
WHT   = RGBColor(0xFF, 0xFF, 0xFF)  # white
DTX   = RGBColor(0x1A, 0x1A, 0x2E)  # dark text
MTX   = RGBColor(0x55, 0x55, 0x65)  # mid text
LTX   = RGBColor(0xB8, 0xD4, 0xEE)  # light text (on dark)
GRN   = RGBColor(0x14, 0x84, 0x40)  # green
RD    = RGBColor(0xC0, 0x22, 0x22)  # red
AMB   = RGBColor(0xB8, 0x78, 0x00)  # amber
LGRN  = RGBColor(0xD5, 0xF0, 0xE0)  # light green (cell bg)
LRD   = RGBColor(0xF8, 0xD8, 0xD8)  # light red (cell bg)
LAMB  = RGBColor(0xFD, 0xF2, 0xD0)  # light amber (cell bg)
LBLU  = RGBColor(0xD6, 0xE9, 0xF5)  # light blue (header bg)
LTEAL = RGBColor(0xD0, 0xEE, 0xF2)  # light teal
LORG  = RGBColor(0xFD, 0xEE, 0xD8)  # light orange (featured row)
ALTW  = RGBColor(0xF5, 0xF8, 0xFD)  # alternate table row
DB2   = RGBColor(0x26, 0x50, 0x8A)  # dark blue 2 (stat boxes)

# ── Helpers ───────────────────────────────────────────────────────────────────

def blank_slide(prs):
    return prs.slides.add_slide(prs.slide_layouts[6])

def bg(slide, rgb):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = rgb

def rect(slide, l, t, w, h, fill=None, line=None):
    shp = slide.shapes.add_shape(1, l, t, w, h)
    if fill:
        shp.fill.solid(); shp.fill.fore_color.rgb = fill
    else:
        shp.fill.background()
    if line:
        shp.line.color.rgb = line
    else:
        shp.line.fill.background()
    return shp

def txt(slide, text, l, t, w, h, sz=16, bold=False, col=DTX,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p = tf.paragraphs[0]; p.alignment = align
    r = p.add_run()
    r.text = text; r.font.size = Pt(sz); r.font.bold = bold
    r.font.italic = italic; r.font.color.rgb = col
    return tb

def txt_lines(slide, lines, l, t, w, h, dsz=14, dcol=DTX, dalign=PP_ALIGN.LEFT):
    """lines: str | (str,bold) | (str,bold,sz) | (str,bold,sz,col) | (str,bold,sz,col,align)"""
    tb = slide.shapes.add_textbox(l, t, w, h)
    tb.word_wrap = True; tf = tb.text_frame; tf.word_wrap = True
    for i, ln in enumerate(lines):
        if isinstance(ln, str):   s,b,sz,c,al = ln,False,dsz,dcol,dalign
        elif len(ln)==2:          s,b=ln; sz,c,al=dsz,dcol,dalign
        elif len(ln)==3:          s,b,sz=ln; c,al=dcol,dalign
        elif len(ln)==4:          s,b,sz,c=ln; al=dalign
        else:                     s,b,sz,c,al=ln
        p = tf.paragraphs[0] if i==0 else tf.add_paragraph()
        p.alignment = al; r = p.add_run()
        r.text=s; r.font.size=Pt(sz); r.font.bold=b; r.font.color.rgb=c
    return tb

def header(slide, title, subtitle=None):
    rect(slide, 0, 0, W, Inches(1.05), fill=DB)
    rect(slide, 0, Inches(1.05), W, Inches(0.06), fill=TEAL)
    txt(slide, title, Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.72),
        sz=26, bold=True, col=WHT)
    if subtitle:
        txt(slide, subtitle, Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.33),
            sz=12, col=LTX)

def cfmt(cell, text, bg_col=WHT, fg_col=DTX, sz=11, bold=False,
         align=PP_ALIGN.CENTER):
    cell.text = text
    cell.fill.solid(); cell.fill.fore_color.rgb = bg_col
    p = cell.text_frame.paragraphs[0]; p.alignment = align
    if p.runs:
        r = p.runs[0]; r.font.size=Pt(sz); r.font.bold=bold; r.font.color.rgb=fg_col

def make_table(slide, headers, rows, l, t, w, h, col_w, hdr_bg=LBLU, fs=11):
    """
    rows: list of lists; each cell: str | (str,bg) | (str,bg,fg) | (str,bg,fg,bold)
    """
    nc=len(headers); nr=len(rows)
    tbl = slide.shapes.add_table(nr+1, nc, l, t, w, h).table
    tot = sum(col_w)
    for i,cw in enumerate(col_w): tbl.columns[i].width = int(w*cw/tot)
    for j,h_ in enumerate(headers):
        cfmt(tbl.cell(0,j), h_, bg_col=hdr_bg, fg_col=DB, sz=fs, bold=True)
    for i,row in enumerate(rows):
        ab = ALTW if i%2==0 else WHT
        for j,cd in enumerate(row):
            if isinstance(cd,str):         cfmt(tbl.cell(i+1,j), cd, bg_col=ab, sz=fs)
            elif len(cd)==2:               cfmt(tbl.cell(i+1,j), cd[0], bg_col=cd[1], sz=fs)
            elif len(cd)==3:               cfmt(tbl.cell(i+1,j), cd[0], bg_col=cd[1], fg_col=cd[2], sz=fs)
            else:                          cfmt(tbl.cell(i+1,j), cd[0], bg_col=cd[1], fg_col=cd[2], sz=fs, bold=cd[3])

def pic(slide, path, l, t, w, h):
    if os.path.exists(path):
        slide.shapes.add_picture(path, l, t, w, h)
    else:
        rect(slide, l, t, w, h, fill=ALTW, line=MTX)
        txt(slide, os.path.basename(path), l, t+h//3, w, h//3,
            sz=9, col=MTX, align=PP_ALIGN.CENTER)

# ── Presentation ──────────────────────────────────────────────────────────────
prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Título
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, DB)
rect(s, 0, Inches(5.4), W, Inches(2.1), fill=RGBColor(0x0F,0x26,0x52))
rect(s, 0, Inches(2.65), W, Inches(0.07), fill=TEAL)
rect(s, 0, Inches(3.25), W, Inches(0.04), fill=ORG)

txt(s, "Dinâmica Molecular de Inibidores de Tripsinas de Spodoptera",
    Inches(0.5), Inches(0.75), Inches(12.3), Inches(1.65),
    sz=34, bold=True, col=WHT, align=PP_ALIGN.CENTER)
txt(s, "Avaliação comparativa: Peptídeo GORE4 vs. SKTI Kunitz  |  100 ns  |  pH 8,2",
    Inches(0.5), Inches(2.45), Inches(12.3), Inches(0.55),
    sz=17, col=LTX, align=PP_ALIGN.CENTER)

stats = [
    ("7",       "Complexos\nsimulados"),
    ("2",       "Séries de\ninibidores"),
    ("100 ns",  "por simulação"),
    ("pH 8,2",  "Intestino médio\nSpodoptera"),
    ("5",       "Modos de\ninibição"),
]
bw = Inches(2.15); bh = Inches(1.55); bt = Inches(4.5)
sl = Inches(0.75); gap = Inches(0.3)
for k,(num,lab) in enumerate(stats):
    bl = sl + k*(bw+gap)
    rect(s, bl, bt, bw, bh, fill=DB2)
    txt(s, num, bl, bt+Inches(0.1), bw, Inches(0.75),
        sz=28, bold=True, col=TEAL, align=PP_ALIGN.CENTER)
    txt(s, lab, bl, bt+Inches(0.8), bw, Inches(0.65),
        sz=12, col=LTX, align=PP_ALIGN.CENTER)

txt(s, "Laboratório de Bioinformática Estrutural  ·  UFV Viçosa-MG  ·  2026",
    Inches(0.5), Inches(6.9), Inches(12.3), Inches(0.4),
    sz=11, col=MTX, align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Sistema Biológico
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Sistema Biológico", "Tripsinas digestivas de Spodoptera × dois inibidores peptídicos")

# left block — Receptores
rect(s, Inches(0.3), Inches(1.25), Inches(5.2), Inches(5.85), fill=ALTW, line=LBLU)
txt(s, "Receptores — Tripsinas de Spodoptera spp.",
    Inches(0.45), Inches(1.35), Inches(4.9), Inches(0.5),
    sz=13, bold=True, col=DB)

cat_headers = ["Isoforma", "His", "Asp", "Ser", "S1"]
cat_rows = [
    [("ACR157",LGRN,DB,True), ("His69",LGRN,GRN,False), ("Asp114",LGRN,GRN,False), ("Ser211",LGRN,GRN,False), ("Ile205",LGRN,GRN,False)],
    [("QCL936",LGRN,DB,True), ("His92",LGRN,GRN,False),  ("Asp142",LGRN,GRN,False), ("Ser247",LGRN,GRN,False), ("Asp241",LGRN,GRN,False)],
    [("XP273", LAMB,DB,True), ("Tyr83",LAMB,AMB,False),  ("Asp132",LAMB,AMB,False), ("Ser234",LAMB,AMB,False), ("Ile229",LAMB,AMB,False)],
    [("XP352", ALTW,DB,True), ("Arg112",ALTW,DTX,False), ("Asp166",ALTW,DTX,False), ("Ser268",ALTW,DTX,False), ("Asp262",ALTW,DTX,False)],
]
make_table(s, cat_headers, cat_rows,
           Inches(0.45), Inches(1.9), Inches(4.9), Inches(1.9),
           [2,1.2,1.2,1.2,1.2], fs=11)

txt_lines(s, [
    ("Características fisiológicas:", True, 12, DB),
    ("• pH 8,2 — intestino médio alcalino de Spodoptera (midgut pH 9,5–11,0)", False, 11, DTX),
    ("• Temperatura: 300 K (fisiológica larval)", False, 11, DTX),
    ("• Ionização: His catalítica → HIE (forma neutra ativa)", False, 11, DTX),
    ("• Íons: 0,10 M KCl (hemolínfa lepidóptero K⁺-dominada)", False, 11, DTX),
    ("• Serino-proteases: tríade His–Asp–Ser + bolsão S1", False, 11, DTX),
], Inches(0.45), Inches(3.9), Inches(4.9), Inches(2.9))

# right block — Ligantes
rect(s, Inches(5.8), Inches(1.25), Inches(7.2), Inches(2.7), fill=LORG, line=RGBColor(0xF0,0xC8,0x90))
txt(s, "Série GORE4 — Pentapeptídeo inibidor",
    Inches(5.95), Inches(1.35), Inches(7.0), Inches(0.45),
    sz=13, bold=True, col=ORG)
txt_lines(s, [
    ("• 5 resíduos  ·  ~43 átomos pesados", False, 12, DTX),
    ("• Todos os resíduos envolvidos na interface", False, 12, DTX),
    ("• Docking proteína–peptídeo (HADDOCK 2.4)", False, 12, DTX),
    ("• HADDOCK scores: −43 a −68 kcal mol⁻¹", False, 12, DTX),
], Inches(5.95), Inches(1.85), Inches(6.9), Inches(1.0))

rect(s, Inches(5.8), Inches(4.2), Inches(7.2), Inches(2.85), fill=LTEAL, line=RGBColor(0x80,0xCC,0xD8))
txt(s, "Série SKTI — Inibidor Kunitz de Soja",
    Inches(5.95), Inches(4.3), Inches(7.0), Inches(0.45),
    sz=13, bold=True, col=TEAL)
txt_lines(s, [
    ("• 177 resíduos  ·  ~1673 átomos pesados", False, 12, DTX),
    ("• Alça reativa SPYRIRF  ·  P1 = Arg", False, 12, DTX),
    ("• Estabilizado por pontes dissulfeto (Kunitz)", False, 12, DTX),
    ("• Docking proteína–proteína (HADDOCK 2.4)", False, 12, DTX),
    ("• HADDOCK scores: −101 a −128 kcal mol⁻¹", False, 12, DTX),
    ("• Inibidor canônico de serino-proteases (referência)", False, 12, DTX),
], Inches(5.95), Inches(4.78), Inches(6.9), Inches(2.1))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Protocolo de Simulação (Metodologia)
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Protocolo de Simulação", "GROMACS 2026 · AMBER99SB-ILDN · TIP3P · Nextflow DSL2")

# Pipeline visual (boxes + arrows)
steps = [
    ("HADDOCK\nDocking", LBLU, DB),
    ("PrepPH\npdb2pqr\nPROPKA", LTEAL, TEAL),
    ("Box + Ions\nKCl 0,10 M\nK⁺/Cl⁻", LGRN, GRN),
    ("NVT\n200 ps\n300 K", LAMB, AMB),
    ("NPT\n500 ps\nBerendsen", LORG, ORG),
    ("Produção\n100 ns\nP-Rahman", LRD, RD),
    ("Análises\nRMSD/RMSF\nTríade", LBLU, DB),
]
bw2 = Inches(1.6); bh2 = Inches(1.5); bt2 = Inches(1.3)
sl2 = Inches(0.35)
for k,(lab,bg2,fc) in enumerate(steps):
    bl2 = sl2 + k*(bw2+Inches(0.15))
    rect(s, bl2, bt2, bw2, bh2, fill=bg2, line=fc)
    txt(s, lab, bl2, bt2+Inches(0.15), bw2, bh2-Inches(0.15),
        sz=11, bold=False, col=fc, align=PP_ALIGN.CENTER)
    if k < len(steps)-1:
        ax = bl2+bw2; ay = bt2+bh2//2-Inches(0.07)
        rect(s, ax, ay, Inches(0.15), Inches(0.14), fill=MTX)
        txt(s, "▶", ax-Inches(0.02), ay-Inches(0.08), Inches(0.2), Inches(0.3),
            sz=10, col=MTX, align=PP_ALIGN.CENTER)

# params table
p_headers = ["Parâmetro", "Valor", "Justificativa"]
p_rows = [
    ["Campo de força",    "AMBER99SB-ILDN",           "Proteínas globulares — padrão literatura"],
    ["Água",             "TIP3P explícita",            "Compatível com AMBER99SB-ILDN"],
    ["Íons",             "K⁺/Cl⁻ 0,10 M",             "Hemolínfa lepidóptero K⁺-dominada"],
    ["Temperatura",      "300 K",                      "Fisiológica larval Spodoptera (25–30°C)"],
    ["NVT",              "200 ps · V-rescale",         "Mínimo para complexos prot–peptídeo"],
    ["NPT equilíbrio",   "500 ps · Berendsen",         "Estabilização de densidade"],
    ["Produção",         "100 ns · Parrinello-Rahman", "Ensemble NPT correto; independência de réplica"],
    ["Passo de tempo",   "2 fs · LINCS (H-bonds)",     "Padrão estável para proteínas"],
    ["Eletrostática",    "PME · rcut = 1,2 nm",        "Tratamento rigoroso de longo alcance"],
    ["GPU offload",      "nb + PME + bonded (CUDA)",   "RTX 5070 Ti 16 GB — servidor Debian"],
]
make_table(s, p_headers, p_rows,
           Inches(0.35), Inches(3.0), Inches(12.6), Inches(4.1),
           [3, 3.5, 6], fs=11)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — GORE4: Parâmetros de DM
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Série GORE4 — Parâmetros de Dinâmica Molecular",
       "4 sistemas · 100 ns · pH 8,2 · 300 K · KCl 0,10 M")

h4 = ["Sistema", "RMSD bb (nm)", "RMSD lig (nm)", "Contatos (<0,4nm)", "H-bonds", "SASA lig (nm²)", "Status"]
r4 = [
    [("QCL936-GORE4 c3",LGRN,DB,True), ("0,124 ± 0,017",LGRN,DTX,False),
     ("0,229 ± 0,042",LGRN,DTX,False), ("340 ± 41",LGRN,DTX,False),
     ("2,75 ± 1,01",LGRN,DTX,False), ("9,04 ± 0,37",LGRN,DTX,False),
     ("✓ Estável",LGRN,GRN,True)],
    [("ACR157-GORE4 c1",LGRN,DB,True), ("0,165 ± 0,016",LGRN,DTX,False),
     ("0,393 ± 0,068",LGRN,DTX,False), ("256 ± 38",LGRN,DTX,False),
     ("1,90 ± 1,38",LGRN,DTX,False), ("9,90 ± 0,47",LGRN,DTX,False),
     ("✓ Estável",LGRN,GRN,True)],
    [("XP273-GORE4 c1",LGRN,DB,True), ("0,248 ± 0,045",LGRN,DTX,False),
     ("0,317 ± 0,068",LGRN,DTX,False), ("260 ± 60",LGRN,DTX,False),
     ("3,19 ± 1,07",LGRN,DTX,False), ("8,14 ± 0,37",LGRN,DTX,False),
     ("✓ Estável",LGRN,GRN,True)],
    [("XP352-GORE4 c4r3",LRD,DB,True), ("0,886 ± 0,440",LRD,RD,True),
     ("0,400 ± 0,067",LRD,DTX,False), ("63 ± 143",LRD,RD,True),
     ("—",LRD,MTX,False), ("—",LRD,MTX,False),
     ("✗ Dissociação",LRD,RD,True)],
]
make_table(s, h4, r4,
           Inches(0.3), Inches(1.25), Inches(12.7), Inches(2.1),
           [2.8,2,2,2.2,1.8,2,2], fs=11)

# diagnosis box
rect(s, Inches(0.3), Inches(3.55), Inches(12.7), Inches(3.65), fill=ALTW, line=LBLU)
txt(s, "Interpretação dos parâmetros",
    Inches(0.5), Inches(3.65), Inches(12.3), Inches(0.4),
    sz=13, bold=True, col=DB)
txt_lines(s, [
    ("Estabilidade (RMSD bb):", True, 12, TEAL),
    ("  QCL936 (0,124) é o mais estável; ACR157 e XP273 apresentam RMSD comparáveis e dentro do critério de estabilidade (<0,3 nm).", False, 11, DTX),
    ("  XP352 destaca-se com RMSD 0,886 nm sem platô — indicativo de dissociação.", False, 11, RD),
    ("", False, 10, DTX),
    ("Interface (contatos e H-bonds):", True, 12, TEAL),
    ("  Todos os sistemas estáveis superam o limiar de 150 contatos; XP352 colapsa (63 ± 143, DP > média).", False, 11, DTX),
    ("  XP273 apresenta o maior número de H-bonds (3,19) e menor SASA do ligante (8,14 nm²) — melhor enterramento.", False, 11, DTX),
    ("", False, 10, DTX),
    ("Critério de dissociação (XP352):", True, 12, RD),
    ("  RMSD crescente sem platô + contatos bimodais (períodos ~ 0) + todas as distâncias catalíticas > 1,0 nm.", False, 11, DTX),
], Inches(0.5), Inches(4.1), Inches(12.3), Inches(3.0))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — GORE4: Painéis de análise
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Série GORE4 — Painéis de Dinâmica Molecular", "Trajetórias 100 ns @ pH 8,2")

panels_g = [
    (os.path.join(IMG,"GORE4","QCL936","QCL936-GORE4.png"),  "QCL936-GORE4 c3  ✓ Estável"),
    (os.path.join(IMG,"GORE4","ACR157","ACR157-GORE4.png"),  "ACR157-GORE4 c1  ✓ Estável"),
    (os.path.join(IMG,"GORE4","XP273","XP273-GORE4.png"),    "XP273-GORE4 c1   ✓ Estável"),
    (os.path.join(IMG,"GORE4","XP352","XP352-GORE4.png"),    "XP352-GORE4 c4r3  ✗ Dissociação"),
]
pw = Inches(6.1); ph = Inches(2.85)
positions = [
    (Inches(0.2),  Inches(1.2)),
    (Inches(6.6),  Inches(1.2)),
    (Inches(0.2),  Inches(4.2)),
    (Inches(6.6),  Inches(4.2)),
]
label_cols = [GRN, GRN, GRN, RD]
for k,((p_img, label),(pl,pt)) in enumerate(zip(panels_g, positions)):
    pic(s, p_img, pl, pt+Inches(0.28), pw, ph-Inches(0.28))
    txt(s, label, pl, pt, pw, Inches(0.28),
        sz=11, bold=True, col=label_cols[k], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — GORE4: Sítio Catalítico
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Série GORE4 — Posicionamento no Sítio Catalítico",
       "Distâncias mínimas (nm) ligante–resíduos catalíticos  ·  limiar de contato: 0,5 nm")

h6 = ["Sistema", "triad_1 (His/Tyr/Arg)", "triad_2 (Asp)", "triad_3 (Ser)", "triad_4 (S1)", "Modo de Inibição"]
r6 = [
    [("QCL936-c3",LGRN,DB,True),
     ("His92  0,242 ✓",LGRN,GRN,True), ("Asp142  1,72 ✗",LRD,RD,False),
     ("Ser247  1,20 ✗",LRD,RD,False),  ("Asp241  0,227 ✓",LGRN,GRN,True),
     ("His + S1",LTEAL,TEAL,True)],
    [("ACR157-c1",LGRN,DB,True),
     ("His69  <0,50 ✓",LGRN,GRN,True), ("Asp114  >1,00 ✗",LRD,RD,False),
     ("Ser211  0,5–0,9 ~",LAMB,AMB,False), ("Ile205  ~0,45 ~",LAMB,AMB,False),
     ("His + S1\n(parcial)",LTEAL,TEAL,False)],
    [("XP273-c1",LGRN,DB,True),
     ("Tyr83  <0,50 ✓",LGRN,GRN,True), ("Asp132  0,55–0,65 ~",LAMB,AMB,False),
     ("Ser234  0,55–0,65 ~",LAMB,AMB,False), ("Ile229  >1,00 ✗",LRD,RD,False),
     ("Tyr periférico",LAMB,AMB,True)],
    [("XP352-c4r3",LRD,DB,True),
     ("Arg112  >1,00 ✗",LRD,RD,False), ("Asp166  >1,00 ✗",LRD,RD,False),
     ("Ser268  >1,00 ✗",LRD,RD,False),  ("Asp262  >1,00 ✗",LRD,RD,False),
     ("— Dissociação",LRD,RD,True)],
]
make_table(s, h6, r6,
           Inches(0.3), Inches(1.25), Inches(12.7), Inches(2.0),
           [2.5,2.4,2.2,2.2,2.2,2.8], fs=11)

# interpretation
rect(s, Inches(0.3), Inches(3.45), Inches(12.7), Inches(3.8), fill=ALTW, line=LBLU)
txt(s, "Interpretação mecanística — Série GORE4",
    Inches(0.5), Inches(3.55), Inches(12.3), Inches(0.45),
    sz=13, bold=True, col=DB)
txt_lines(s, [
    ("Modo His+S1 (QCL936 e ACR157):", True, 12, TEAL),
    ("  O GORE4 ancora na His da tríade e no bolsão S1, sem contato com o díade Asp–Ser. Bloqueia o reconhecimento de substrato (S1)", False, 11, DTX),
    ("  e interfere perifericamente na tríade catalítica. QCL936 apresenta interface mais compacta (340 vs 256 contatos).", False, 11, DTX),
    ("", False, 10, DTX),
    ("Modo Tyr periférico (XP273):", True, 12, AMB),
    ("  Ancoragem exclusivamente em Tyr83 (resíduo periférico atípico). Asp132/Ser234 borderline; Ile229/S1 distante.", False, 11, DTX),
    ("  Mecanismo por oclusão estérica da entrada do sítio ou estabilização de conformação cataliticamente inativa.", False, 11, DTX),
    ("  SASA mais baixa (8,14 nm²) e maior H-bonds (3,19) — enterramento eficiente mesmo sem contato S1.", False, 11, DTX),
    ("", False, 10, DTX),
    ("Dissociação (XP352-rank3):", True, 12, RD),
    ("  Pose inicial de baixa qualidade (rank 3, cluster 4). RMSD 0,886 nm; contatos SD > média — instabilidade estrutural.", False, 11, DTX),
], Inches(0.5), Inches(4.05), Inches(12.3), Inches(3.1))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — SKTI: Parâmetros de DM
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Série SKTI — Parâmetros de Dinâmica Molecular",
       "3 sistemas · SKTI 177 aa · alça reativa SPYRIRF · P1=Arg · 100 ns @ pH 8,2")

h7 = ["Sistema", "RMSD bb (nm)", "RMSD lig (nm)", "Rg (nm)", "Contatos (<0,4nm)", "H-bonds", "SASA lig (nm²)", "Status"]
r7 = [
    [("ACR157-SKTI c2",LGRN,DB,True), ("0,281 ± 0,027",LGRN,DTX,False),
     ("0,206 ± 0,017",LGRN,GRN,True), ("2,225 ± 0,009",LGRN,DTX,False),
     ("1018 ± 118",LGRN,GRN,True), ("13,524 ± 2,911",LGRN,GRN,True),
     ("101,97 ± 1,02",LGRN,DTX,False), ("✓ Estável",LGRN,GRN,True)],
    [("QCL936-SKTI c2",LORG,DB,True), ("0,370 ± 0,054",LORG,AMB,False),
     ("0,209 ± 0,025",LORG,DTX,False), ("2,296 ± 0,015",LORG,DTX,False),
     ("720 ± 77",LORG,DTX,False), ("10,410 ± 2,022",LORG,DTX,False),
     ("102,53 ± 2,09",LORG,DTX,False), ("✓ Estável",LGRN,GRN,True)],
    [("XP273-SKTI c2",LGRN,DB,True), ("0,288 ± 0,042",LGRN,DTX,False),
     ("0,224 ± 0,017",LGRN,DTX,False), ("2,325 ± 0,015",LGRN,DTX,False),
     ("596 ± 88",LGRN,DTX,False), ("8,817 ± 3,496",LGRN,DTX,False),
     ("101,31 ± 2,01",LGRN,DTX,False), ("✓ Estável",LGRN,GRN,True)],
]
make_table(s, h7, r7,
           Inches(0.3), Inches(1.25), Inches(12.7), Inches(1.85),
           [2.8,2,2,1.8,2.2,2.2,2.2,1.8], fs=11)

# comparison with GORE4
rect(s, Inches(0.3), Inches(3.25), Inches(12.7), Inches(4.0), fill=ALTW, line=LBLU)
txt(s, "SKTI vs GORE4 — comparação quantitativa por receptor",
    Inches(0.5), Inches(3.35), Inches(12.3), Inches(0.4),
    sz=13, bold=True, col=DB)

h_cmp = ["Receptor", "Contatos GORE4", "Contatos SKTI", "Razão", "H-bonds GORE4", "H-bonds SKTI", "Razão"]
r_cmp = [
    ["ACR157", "256 ± 38", ("1018 ± 118",LGRN), ("3,98×",LGRN,GRN,True), "1,90 ± 1,38", ("13,524 ± 2,911",LGRN), ("7,1×",LGRN,GRN,True)],
    ["QCL936", "340 ± 41", ("720 ± 77",LAMB),   ("2,12×",LAMB,AMB,False), "2,75 ± 1,01", ("10,410 ± 2,022",LAMB), ("3,8×",LAMB,AMB,False)],
    ["XP273",  "260 ± 60", ("596 ± 88",LTEAL),  ("2,29×",LTEAL,TEAL,False), "3,19 ± 1,07", ("8,817 ± 3,496",LTEAL), ("2,8×",LTEAL,TEAL,False)],
]
make_table(s, h_cmp, r_cmp,
           Inches(0.5), Inches(3.85), Inches(12.3), Inches(1.45),
           [2.2,2,2,1.5,2.2,2.5,1.5], fs=11)

txt(s, "A série SKTI supera a GORE4 em contatos (2,1–4,0×) e H-bonds (2,8–7,1×) em todos os receptores — coerente com a natureza proteica e pré-organizada do inibidor Kunitz.",
    Inches(0.5), Inches(5.45), Inches(12.3), Inches(0.5),
    sz=12, col=TEAL, bold=True)
txt_lines(s, [
    ("Todos os três complexos SKTI são estáveis (RMSD lig < 0,25 nm), ao contrário do GORE4 em ACR157 (0,393 nm).", False, 11, DTX),
    ("QCL936-SKTI apresenta RMSD do backbone mais elevado (0,370 nm) — maior rearranjo conformacional induzido pelo SKTI volumoso.", False, 11, MTX),
], Inches(0.5), Inches(6.0), Inches(12.3), Inches(0.9))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — SKTI: Painéis
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Série SKTI — Painéis de Dinâmica Molecular", "Trajetórias 100 ns @ pH 8,2")

panels_s = [
    (os.path.join(IMG,"SKTI","ACR157-SKTI-painel_completo.png"), "ACR157-SKTI c2  ·  Kunitz canônico 4/4"),
    (os.path.join(IMG,"SKTI","QCL936-SKTI-painel_completo.png"), "QCL936-SKTI c2  ·  His+Ser+S1 3/4  [novo]"),
    (os.path.join(IMG,"SKTI","XP273-SKTI-painel_completo.png"),  "XP273-SKTI c2   ·  Tyr+Asp+S1 3/4"),
]
pw3 = Inches(4.15); ph3 = Inches(5.8)
pl3 = [Inches(0.2), Inches(4.6), Inches(9.0)]
label_c3 = [GRN, ORG, TEAL]
for k,(p_img,label) in enumerate(panels_s):
    pic(s, p_img, pl3[k], Inches(1.55), pw3, ph3)
    txt(s, label, pl3[k], Inches(1.2), pw3, Inches(0.35),
        sz=11, bold=True, col=label_c3[k], align=PP_ALIGN.CENTER)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — QCL936-SKTI: Descoberta do re-run
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
rect(s, 0, 0, W, Inches(1.05), fill=ORG)
rect(s, 0, Inches(1.05), W, Inches(0.06), fill=TEAL)
txt(s, "QCL936-SKTI: Descoberta Crítica do Re-run",
    Inches(0.4), Inches(0.1), Inches(12.5), Inches(0.72),
    sz=26, bold=True, col=WHT)
txt(s, "Correção dos resíduos catalíticos de Asp114/Ser211 → Asp142/Ser247 revela modo His+Ser+S1",
    Inches(0.4), Inches(0.75), Inches(12.5), Inches(0.33), sz=12, col=WHT)

# Before / After boxes
rect(s, Inches(0.3), Inches(1.25), Inches(5.9), Inches(3.3), fill=LRD, line=RD)
txt(s, "ANTES — Resíduos Incorretos (Asp114/Ser211)",
    Inches(0.45), Inches(1.35), Inches(5.7), Inches(0.45),
    sz=13, bold=True, col=RD)
h_bef = ["Resíduo", "Distância (nm)", "Engajamento"]
r_bef = [
    [("His92",LGRN), ("0,179",LGRN,GRN,True), ("✓ Contato",LGRN,GRN,True)],
    [("Asp114 ✗",LRD), ("1,38",LRD,RD,False), ("✗ Distante",LRD,RD,True)],
    [("Ser211 ✗",LRD), ("1,08",LRD,RD,False), ("✗ Distante",LRD,RD,True)],
    [("Asp241/S1",LGRN), ("0,178",LGRN,GRN,True), ("✓ Contato",LGRN,GRN,True)],
]
make_table(s, h_bef, r_bef,
           Inches(0.45), Inches(1.85), Inches(5.6), Inches(1.65),
           [2,2,2.5], fs=11)
txt(s, "Interpretação incorreta: modo His+S1 (2/4)",
    Inches(0.45), Inches(3.55), Inches(5.6), Inches(0.45),
    sz=12, bold=True, col=RD, align=PP_ALIGN.CENTER)
txt(s, "Ser nucleofílica interpretada como DISTANTE",
    Inches(0.45), Inches(3.95), Inches(5.6), Inches(0.45),
    sz=11, col=RD, align=PP_ALIGN.CENTER)

rect(s, Inches(6.5), Inches(1.25), Inches(6.5), Inches(3.3), fill=LGRN, line=GRN)
txt(s, "DEPOIS — Resíduos Corretos (Asp142/Ser247)",
    Inches(6.65), Inches(1.35), Inches(6.3), Inches(0.45),
    sz=13, bold=True, col=GRN)
h_aft = ["Resíduo", "Distância (nm)", "Engajamento"]
r_aft = [
    [("His92",LGRN), ("0,179",LGRN,GRN,True), ("✓ Contato",LGRN,GRN,True)],
    [("Asp142",LRD), ("~0,73",LRD,RD,False), ("✗ Distante",LRD,RD,True)],
    [("Ser247 NOVO",LGRN), ("0,176",LGRN,GRN,True), ("✓ Contato",LGRN,GRN,True)],
    [("Asp241/S1",LGRN), ("0,178",LGRN,GRN,True), ("✓ Contato",LGRN,GRN,True)],
]
make_table(s, h_aft, r_aft,
           Inches(6.65), Inches(1.85), Inches(6.2), Inches(1.65),
           [2,2,2.5], fs=11)
txt(s, "Modo correto: His+Ser+S1 (3/4)  — Ser247 ENGAJADA",
    Inches(6.65), Inches(3.55), Inches(6.2), Inches(0.45),
    sz=12, bold=True, col=GRN, align=PP_ALIGN.CENTER)
txt(s, "Asp142 permanece livre — etapa de regeneração da His não bloqueada",
    Inches(6.65), Inches(3.95), Inches(6.2), Inches(0.45),
    sz=11, col=AMB, align=PP_ALIGN.CENTER)

# Implication box
rect(s, Inches(0.3), Inches(4.75), Inches(12.7), Inches(2.55), fill=LORG, line=ORG)
txt(s, "Implicação Mecanística",
    Inches(0.5), Inches(4.85), Inches(12.3), Inches(0.45),
    sz=13, bold=True, col=ORG)
txt_lines(s, [
    ("Com os resíduos corretos, o SKTI engaja simultaneamente His92 e Ser247 — os dois componentes nucleofílicos da tríade — mais o bolsão S1 (Asp241).", False, 11, DTX),
    ("O Asp142 (recarregador da His na catálise) permanece livre: a etapa de regeneração catalítica não está bloqueada estruturalmente,", False, 11, DTX),
    ("o que diferencia este modo (3/4) do mecanismo Kunitz canônico do ACR157 (4/4), onde todos os componentes são simultaneamente bloqueados.", False, 11, DTX),
    ("", False, 10, DTX),
    ("Este resultado eleva QCL936-SKTI de 2/4 para 3/4 de engajamento, revisando a hierarquia de modos de inibição da série SKTI.", False, 11, ORG, PP_ALIGN.LEFT),
], Inches(0.5), Inches(5.35), Inches(12.3), Inches(1.8))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — SKTI: Sítio Catalítico
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Série SKTI — Posicionamento no Sítio Catalítico",
       "Distâncias mínimas (nm) SKTI–resíduos catalíticos  ·  limiar de contato: 0,5 nm")

h10 = ["Sistema", "triad_1 (His/Tyr)", "triad_2 (Asp)", "triad_3 (Ser)", "triad_4 (S1)", "Modo"]
r10 = [
    [("ACR157-c2",LGRN,DB,True),
     ("His69  0,224 ✓",LGRN,GRN,True), ("Asp114  0,312 ✓",LGRN,GRN,True),
     ("Ser211  0,221 ✓",LGRN,GRN,True), ("Ile205  0,180 ✓",LGRN,GRN,True),
     ("Kunitz canônico\n4/4 ✓✓✓✓",LGRN,GRN,True)],
    [("QCL936-c2",LORG,DB,True),
     ("His92  0,179 ✓",LGRN,GRN,True), ("Asp142  ~0,73 ✗",LRD,RD,False),
     ("Ser247  0,176 ✓",LGRN,GRN,True), ("Asp241  0,178 ✓",LGRN,GRN,True),
     ("His+Ser+S1\n3/4 ✓✗✓✓",LORG,ORG,True)],
    [("XP273-c2",LGRN,DB,True),
     ("Tyr83  0,098 ✓",LGRN,GRN,True), ("Asp132  0,119 ✓",LGRN,GRN,True),
     ("Ser234  0,718 ✗",LRD,RD,False), ("Ile229  0,118 ✓",LGRN,GRN,True),
     ("Tyr+Asp+S1\n3/4 ✓✓✗✓",LTEAL,TEAL,True)],
]
make_table(s, h10, r10,
           Inches(0.3), Inches(1.25), Inches(12.7), Inches(2.15),
           [2.6,2.5,2.5,2.5,2.5,2.9], fs=12)

# interpretation
rect(s, Inches(0.3), Inches(3.6), Inches(12.7), Inches(3.7), fill=ALTW, line=LBLU)
txt(s, "Interpretação comparativa dos modos SKTI",
    Inches(0.5), Inches(3.7), Inches(12.3), Inches(0.4),
    sz=13, bold=True, col=DB)
txt_lines(s, [
    ("ACR157-SKTI — Kunitz canônico completo (4/4):", True, 12, GRN),
    ("  His69 (0,224) + Asp114 (0,312) + Ser211 (0,221) + Ile205/S1 (0,180) — todos em contato. P1=Arg insere-se diretamente no núcleo catalítico.", False, 11, DTX),
    ("  O bloqueio simultâneo de His–Asp–Ser impede o complexo tetraédrico de transição; S1 bloqueado simultaneamente.", False, 11, DTX),
    ("", False, 10, DTX),
    ("QCL936-SKTI — His + Ser + S1 (3/4), modo único:", True, 12, ORG),
    ("  His92 e Ser247 engajadas (ambas nucleofílicas na catálise), mas Asp142 (recarregador da His) fica livre.", False, 11, DTX),
    ("  Forte atração eletrostática Asp241(S1)–P1=Arg direciona SPYRIRF para His+Ser sem alcançar Asp142.", False, 11, DTX),
    ("", False, 10, DTX),
    ("XP273-SKTI — Tyr + Asp + S1 (3/4):", True, 12, TEAL),
    ("  Tyr83 (0,098) + Asp132 (0,119) + Ile229/S1 (0,118) engajados; Ser234 (0,718) livre.", False, 11, DTX),
    ("  Presença de Tyr83 no lugar da His canônica reorienta SPYRIRF — Ser não alcançada geometricamente.", False, 11, DTX),
], Inches(0.5), Inches(4.15), Inches(12.3), Inches(3.0))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — Comparação GORE4 vs SKTI + 5 Modos
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Síntese Comparativa — GORE4 vs SKTI e 5 Modos de Inibição", "")

# Tabela 1 completa compacta
rect(s, Inches(0.3), Inches(1.2), Inches(8.0), Inches(0.35), fill=DB)
txt(s, "Tabela 1 — Parâmetros de DM (médias ± DP, 100 ns)",
    Inches(0.4), Inches(1.22), Inches(7.8), Inches(0.3),
    sz=11, bold=True, col=WHT)
h11a = ["Sistema", "RMSD bb", "RMSD lig", "Contatos", "H-bonds", "Status"]
r11a = [
    [("QCL936-GORE4",LGRN,DB,True), "0,124±0,017", "0,229±0,042", "340±41", "2,75±1,01", ("✓ Estável",LGRN,GRN,True)],
    [("ACR157-GORE4",LGRN,DB,True), "0,165±0,016", "0,393±0,068", "256±38", "1,90±1,38", ("✓ Estável",LGRN,GRN,True)],
    [("XP273-GORE4", LGRN,DB,True), "0,248±0,045", "0,317±0,068", "260±60", "3,19±1,07", ("✓ Estável",LGRN,GRN,True)],
    [("XP352-GORE4", LRD, DB,True), "0,886±0,440", "0,400±0,067", "63±143", "—",         ("✗ Dissoc.",LRD,RD,True)],
    [("ACR157-SKTI", LGRN,DB,True), "0,281±0,027", "0,206±0,017", "1018±118","13,524±2,911",("✓ Estável",LGRN,GRN,True)],
    [("QCL936-SKTI", LORG,DB,True), "0,370±0,054", "0,209±0,025", "720±77",  "10,410±2,022",("✓ Estável",LGRN,GRN,True)],
    [("XP273-SKTI",  LGRN,DB,True), "0,288±0,042", "0,224±0,017", "596±88",  "8,817±3,496", ("✓ Estável",LGRN,GRN,True)],
]
make_table(s, h11a, r11a,
           Inches(0.3), Inches(1.55), Inches(8.0), Inches(3.3),
           [2.6,1.8,1.8,1.8,2.2,1.5], fs=10)

# 5 Modos box
rect(s, Inches(8.5), Inches(1.2), Inches(4.6), Inches(6.1), fill=ALTW, line=LBLU)
txt(s, "5 Modos de Inibição Identificados",
    Inches(8.65), Inches(1.3), Inches(4.3), Inches(0.45),
    sz=12, bold=True, col=DB, align=PP_ALIGN.CENTER)
modos = [
    (LGRN, GRN, "1. Kunitz canônico (4/4)", "ACR157-SKTI", "His+Asp+Ser+S1"),
    (LORG, ORG, "2. His+Ser+S1 (3/4)",      "QCL936-SKTI", "His+Ser+S1·Asp livre"),
    (LTEAL,TEAL,"3. Tyr+Asp+S1 (3/4)",      "XP273-SKTI",  "Tyr+Asp+S1·Ser livre"),
    (LAMB, AMB, "4. His+S1 (2/4)",           "QCL936/ACR157-GORE4","His+S1·Asp-Ser livres"),
    (LBLU, DB,  "5. Periférico (1/4)",       "XP273-GORE4", "Tyr83 apenas"),
]
for k,(bg_m,fg_m,titulo,sistema,desc) in enumerate(modos):
    my = Inches(1.85) + k*Inches(1.1)
    rect(s, Inches(8.65), my, Inches(4.2), Inches(0.95), fill=bg_m, line=fg_m)
    txt(s, titulo, Inches(8.75), my+Inches(0.05), Inches(4.0), Inches(0.38),
        sz=11, bold=True, col=fg_m)
    txt(s, f"  {sistema}", Inches(8.75), my+Inches(0.38), Inches(4.0), Inches(0.28),
        sz=10, col=DTX)
    txt(s, f"  {desc}", Inches(8.75), my+Inches(0.62), Inches(4.0), Inches(0.28),
        sz=10, col=MTX, italic=True)

# summary line
txt(s, "SKTI supera GORE4: contatos 2,1–4,0×  ·  H-bonds 2,8–7,1×  ·  RMSD lig inferior em todos os receptores",
    Inches(0.3), Inches(5.0), Inches(7.9), Inches(0.45),
    sz=11, bold=True, col=TEAL, align=PP_ALIGN.CENTER)
# Tabela 2 compacta
rect(s, Inches(0.3), Inches(5.55), Inches(8.0), Inches(0.3), fill=DB)
txt(s, "Tabela 2 — Sítio catalítico (distâncias médias, nm)",
    Inches(0.4), Inches(5.57), Inches(7.8), Inches(0.26),
    sz=10, bold=True, col=WHT)
h11b = ["Sistema", "triad_1", "triad_2(Asp)", "triad_3(Ser)", "triad_4(S1)", "Modo"]
r11b = [
    ["QCL936-GORE4", ("0,242✓",LGRN,GRN,False), ("1,72✗",LRD,RD,False), ("1,20✗",LRD,RD,False), ("0,227✓",LGRN,GRN,False), "His+S1"],
    ["ACR157-GORE4", ("<0,50✓",LGRN,GRN,False), (">1,00✗",LRD,RD,False), ("~✗",LRD,RD,False),    ("~0,45~",LAMB,AMB,False), "His+S1 parcial"],
    ["XP273-GORE4",  ("<0,50✓",LGRN,GRN,False), ("0,55–0,65~",LAMB,AMB,False), ("0,55–0,65~",LAMB,AMB,False), (">1,00✗",LRD,RD,False), "Tyr periférico"],
    ["ACR157-SKTI",  ("0,224✓",LGRN,GRN,True),  ("0,312✓",LGRN,GRN,True), ("0,221✓",LGRN,GRN,True), ("0,180✓",LGRN,GRN,True), ("Kunitz 4/4",LGRN,GRN,True)],
    ["QCL936-SKTI",  ("0,179✓",LGRN,GRN,True),  ("~0,73✗",LRD,RD,False),  ("0,176✓",LGRN,GRN,True), ("0,178✓",LGRN,GRN,True), ("H+Ser+S1 3/4",LORG,ORG,True)],
    ["XP273-SKTI",   ("0,098✓",LGRN,GRN,True),  ("0,119✓",LGRN,GRN,True), ("0,718✗",LRD,RD,False), ("0,118✓",LGRN,GRN,True), ("Tyr+D+S1 3/4",LTEAL,TEAL,True)],
]
make_table(s, h11b, r11b,
           Inches(0.3), Inches(5.85), Inches(8.0), Inches(1.5),
           [2.2,1.6,1.8,1.8,1.6,2.2], fs=9)

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — Hierarquia e Ranking
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Hierarquia e Ranking — SKTI vs GORE4",
       "Validação cruzada: escores HADDOCK × parâmetros de DM × profundidade de engajamento catalítico")

# Ranking SKTI
rect(s, Inches(0.3), Inches(1.2), Inches(6.1), Inches(0.35), fill=DB)
txt(s, "Ranking — Série SKTI", Inches(0.4), Inches(1.22), Inches(5.9), Inches(0.3),
    sz=12, bold=True, col=WHT)
rank_s = [
    ("1º", "ACR157-SKTI", "4/4 — Kunitz canônico", "1018", "13,524", "-128,3", GRN,  LGRN),
    ("2º", "QCL936-SKTI", "3/4 — His+Ser+S1",       "720",  "10,410", "-113,8", ORG,  LORG),
    ("3º", "XP273-SKTI",  "3/4 — Tyr+Asp+S1",       "596",  "8,817",  "-101,4", TEAL, LTEAL),
]
for k,(pos,sis,modo,cont,hb,hadd,fc,bg_r) in enumerate(rank_s):
    ry = Inches(1.6) + k*Inches(1.4)
    rect(s, Inches(0.3), ry, Inches(6.1), Inches(1.25), fill=bg_r, line=fc)
    txt(s, pos, Inches(0.4), ry+Inches(0.1), Inches(0.6), Inches(0.8),
        sz=28, bold=True, col=fc, align=PP_ALIGN.CENTER)
    txt(s, sis, Inches(1.05), ry+Inches(0.08), Inches(3.5), Inches(0.45),
        sz=13, bold=True, col=DTX)
    txt(s, modo, Inches(1.05), ry+Inches(0.52), Inches(3.5), Inches(0.38),
        sz=11, col=fc, bold=True)
    txt(s, f"Contatos: {cont}  ·  H-bonds: {hb}  ·  HADDOCK: {hadd} kcal/mol",
        Inches(1.05), ry+Inches(0.88), Inches(5.2), Inches(0.33),
        sz=10, col=MTX)

# Ranking GORE4
rect(s, Inches(0.3), Inches(5.85), Inches(6.1), Inches(0.3), fill=DB)
txt(s, "Série GORE4 (estáveis)", Inches(0.4), Inches(5.87), Inches(5.9), Inches(0.26),
    sz=11, bold=True, col=WHT)
rank_g = [("QCL936","His+S1","340","2,75","-45,7"),
          ("ACR157","His+S1 parcial","256","1,90","-46,1"),
          ("XP273","Tyr periférico","260","3,19","-43,1")]
make_table(s, ["Receptor","Modo","Contatos","H-bonds","HADDOCK"],
           [[r[0],r[1],r[2],r[3],r[4]] for r in rank_g],
           Inches(0.3), Inches(6.15), Inches(6.1), Inches(1.1),
           [1.8,2.5,1.5,1.6,1.8], fs=10)

# validation and insight box
rect(s, Inches(6.7), Inches(1.2), Inches(6.3), Inches(6.05), fill=ALTW, line=LBLU)
txt(s, "Validação Cruzada e Insights",
    Inches(6.9), Inches(1.3), Inches(6.0), Inches(0.4),
    sz=13, bold=True, col=DB)
txt_lines(s, [
    ("Hierarquia SKTI: ACR157 > QCL936 > XP273", True, 12, TEAL),
    ("  Consistente com escores HADDOCK (−128 > −114 > −101).", False, 11, DTX),
    ("  Validação retroativa: docking previu corretamente a qualidade da interface.", False, 11, DTX),
    ("", False, 10, DTX),
    ("Hierarquia GORE4: QCL936 > XP273 ≈ ACR157", True, 12, ORG),
    ("  QCL936 tem mais contatos (340) e RMSD lig menor (0,229).", False, 11, DTX),
    ("  ACR157 tem interface menos compacta mas mecanismo dual parcial.", False, 11, DTX),
    ("", False, 10, DTX),
    ("SKTI >> GORE4 em todos os receptores:", True, 12, GRN),
    ("  Contatos: 2,1–4,0× maiores  ·  H-bonds: 2,8–7,1× maiores", False, 11, DTX),
    ("  Engajamento catalítico: GORE4 máx 2/4 vs SKTI máx 4/4", False, 11, DTX),
    ("", False, 10, DTX),
    ("Papel do bolsão S1 (Asp241 QCL936 vs Ile205 ACR157):", True, 12, AMB),
    ("  Asp241⁻ (QCL936) atrai P1=Arg fortemente → orienta SPYRIRF", False, 11, DTX),
    ("    para His+Ser+S1, sem alcançar Asp142.", False, 11, DTX),
    ("  Ile205 neutro (ACR157) permite orientação canônica completa.", False, 11, DTX),
    ("", False, 10, DTX),
    ("Tyr83 em XP273 — topografia atípica:", True, 12, TEAL),
    ("  Substituição His → Tyr altera geometria de entrada do sítio.", False, 11, DTX),
    ("  Reorienta SPYRIRF: Ser234 não é alcançada em nenhuma das", False, 11, DTX),
    ("  duas séries de inibidores (GORE4 borderline; SKTI distante).", False, 11, DTX),
], Inches(6.9), Inches(1.75), Inches(5.9), Inches(5.3))

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 13 — Discussão: Seletividade e Implicações
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Discussão — Seletividade por Isoforma e Implicações para Bioinseticidas",
       "Diferentes geometrias do sítio ativo determinam o modo de engajamento do inibidor")

# Three columns
cols = [
    ("ACR157", GRN, LGRN, [
        "Sítio ativo compatível com\nmecanismo Kunitz canônico",
        "",
        "SKTI: bloqueio integral 4/4\n→ máxima potência de inibição",
        "",
        "GORE4: His+S1 parcial\n→ eficiência intermediária",
        "",
        "Ile205/S1 neutro permite\norientação canônica SPYRIRF",
        "",
        "Melhor alvo para design de\ninibidores baseados em Kunitz",
    ]),
    ("QCL936", ORG, LORG, [
        "Asp241 carregado (S1−)\natrai P1=Arg fortemente",
        "",
        "SKTI: modo His+Ser+S1 (3/4)\n→ Asp142 livre (novo resultado)",
        "",
        "GORE4: His+S1 (2/4)\n→ mesmo padrão S1-dominante",
        "",
        "Topografia S1 direciona\ninibidor para His+S1/Ser+S1",
        "",
        "Alvo viável — Asp241\npode ser explorado no design",
    ]),
    ("XP273", TEAL, LTEAL, [
        "Tyr83 no lugar da His\ncalítica — topografia única",
        "",
        "SKTI: Tyr+Asp+S1 (3/4)\n→ Ser234 não alcançada",
        "",
        "GORE4: Tyr periférico (1/4)\n→ menor engajamento catalítico",
        "",
        "Menor complementaridade\nde interface (596 contatos)",
        "",
        "Requer design específico para\ngeometria Tyr83",
    ]),
]
cw3 = Inches(4.0); ct = Inches(1.2); ch = Inches(6.1)
cl_list = [Inches(0.4), Inches(4.6), Inches(8.8)]
for k,(titulo,fc,bg_c,lines) in enumerate(cols):
    cl = cl_list[k]
    rect(s, cl, ct, cw3, Inches(0.5), fill=fc)
    txt(s, titulo, cl, ct+Inches(0.05), cw3, Inches(0.4),
        sz=18, bold=True, col=WHT, align=PP_ALIGN.CENTER)
    rect(s, cl, ct+Inches(0.5), cw3, ch-Inches(0.5), fill=bg_c, line=fc)
    y_off = 0
    for line in lines:
        if line == "":
            y_off += 0.22
            continue
        txt(s, line, cl+Inches(0.1), ct+Inches(0.6)+Inches(y_off),
            cw3-Inches(0.2), Inches(0.55),
            sz=11, col=DTX if not line.startswith("→") else fc)
        y_off += 0.58

# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 14 — Conclusões e Perspectivas
# ══════════════════════════════════════════════════════════════════════════════
s = blank_slide(prs)
bg(s, WHT)
header(s, "Conclusões e Perspectivas",
       "Dinâmica Molecular 100 ns — GORE4 e SKTI vs Tripsinas de Spodoptera")

rect(s, Inches(0.3), Inches(1.2), Inches(7.8), Inches(5.9), fill=ALTW, line=LBLU)
txt(s, "Principais Conclusões", Inches(0.5), Inches(1.3), Inches(7.6), Inches(0.4),
    sz=14, bold=True, col=DB)
conclusions = [
    (GRN,  "1.", "Série SKTI supera GORE4 em todos os receptores: contatos 2,1–4,0×, H-bonds 2,8–7,1×,\n    interface mais estável e engajamento catalítico mais profundo."),
    (GRN,  "2.", "ACR157-SKTI: mecanismo Kunitz canônico completo (4/4) — único complexo com bloqueio\n    simultâneo de His–Asp–Ser e S1."),
    (ORG,  "3.", "QCL936-SKTI: modo His+Ser+S1 (3/4) — descoberta do re-run com resíduos corretos\n    (Asp142/Ser247). Ser247 engajada; Asp142 livre."),
    (TEAL, "4.", "XP273-SKTI: modo Tyr+Asp+S1 (3/4) — geometria atípica Tyr83 impede contato com Ser234."),
    (AMB,  "5.", "Série GORE4: QCL936/ACR157 operam por modo His+S1 (bloqueio de reconhecimento de\n    substrato); XP273-GORE4 por ancoragem periférica exclusiva em Tyr83."),
    (RD,   "6.", "XP352-GORE4 c4r3: dissociação confirmada — pose inicial de baixa qualidade (rank3)."),
    (DB,   "7.", "Hierarquia consistente com HADDOCK em ambas as séries: ACR157 > QCL936 > XP273\n    (valida docking como preditor de qualidade de interface)."),
    (TEAL, "8.", "5 modos de inibição identificados, ordenados por completude: canônico > His+Ser+S1 >\n    Tyr+Asp+S1 > His+S1 > periférico."),
]
for k,(col,num,text) in enumerate(conclusions):
    ky = Inches(1.8) + k*Inches(0.6)
    txt(s, num, Inches(0.45), ky, Inches(0.4), Inches(0.55),
        sz=12, bold=True, col=col)
    txt(s, text, Inches(0.85), ky, Inches(7.1), Inches(0.55),
        sz=11, col=DTX)

# Perspectives box
rect(s, Inches(8.4), Inches(1.2), Inches(4.7), Inches(5.9), fill=LORG, line=ORG)
txt(s, "Perspectivas", Inches(8.6), Inches(1.3), Inches(4.4), Inches(0.4),
    sz=14, bold=True, col=ORG)
txt_lines(s, [
    ("Análises complementares:", True, 12, DB),
    ("• Contact maps (MDAnalysis)\n  heatmap resíduo × resíduo", False, 11, DTX),
    ("• Interaction fingerprints (ProLIF)\n  persistência temporal de interações", False, 11, DTX),
    ("• Snapshots estruturais (PyMOL batch)\n  visualização do modo de ligação", False, 11, DTX),
    ("", False, 10, DTX),
    ("Sistemas pendentes:", True, 12, DB),
    ("• DN773-GORE4 e DN1441-GORE4\n  PDBs disponíveis para simulação", False, 11, DTX),
    ("• XP352-SKTI\n  complexo ainda não simulado", False, 11, DTX),
    ("• XP352-GORE4 rank1\n  re-run com cluster de qualidade", False, 11, DTX),
    ("", False, 10, DTX),
    ("Design de inibidores:", True, 12, DB),
    ("• ACR157: alvo prioritário — aceita\n  mecanismo canônico Kunitz", False, 11, DTX),
    ("• QCL936: design para Asp142\n  pode elevar 3/4 → 4/4", False, 11, DTX),
    ("• XP273: requer design adaptado\n  à topografia Tyr83", False, 11, DTX),
], Inches(8.6), Inches(1.8), Inches(4.4), Inches(5.0))

# Footer
rect(s, 0, Inches(7.2), W, Inches(0.3), fill=DB)
txt(s, "Laboratório de Bioinformática Estrutural · UFV Viçosa-MG · 2026  |  Pipeline Nextflow DSL2 + GROMACS 2026 + HADDOCK 2.4",
    Inches(0.3), Inches(7.22), Inches(13.0), Inches(0.26),
    sz=9, col=LTX, align=PP_ALIGN.CENTER)

# ── Salvar ────────────────────────────────────────────────────────────────────
out = os.path.join(BASE, "apresentacao_md_spodoptera.pptx")
prs.save(out)
print(f"Salvo: {out}")
print(f"Total de slides: {len(prs.slides)}")
